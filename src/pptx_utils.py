from typing import List, Tuple
from pptx import Presentation


def extract_pptx_text_with_slides(path: str) -> List[Tuple[int, str]]:
    # Εξάγει το κείμενο από αρχεία PowerPoint.
    # Επιστρέφει μια λίστα με ζεύγη (αριθμός διαφάνειας, κείμενο),
    # επιτρέποντας την ακριβή αναφορά στην πηγή κατά την αναζήτηση.
    
    # Φορτώνω το αρχείο παρουσίασης χρησιμοποιώντας τη βιβλιοθήκη python-pptx.
    prs = Presentation(path)
    results: List[Tuple[int, str]] = []

    # Διασχίζω όλες τις διαφάνειες της παρουσίασης, ξεκινώντας την αρίθμηση από το 1.
    for i, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        
        # Ελέγχω κάθε αντικείμενο (shape) μέσα στη διαφάνεια για να βρω κείμενο.
        for shape in slide.shapes:
            # Αν το αντικείμενο δεν έχει πλαίσιο κειμένου (π.χ. είναι απλή εικόνα), το αγνοώ.
            if not hasattr(shape, "has_text_frame"):
                continue
            
            if shape.has_text_frame:
                text_runs: List[str] = []
                # Διαβάζω το κείμενο παράγραφο προς παράγραφο για να διατηρήσω τη δομή.
                for paragraph in shape.text_frame.paragraphs:
                    # Ενώνω τα επιμέρους τμήματα (runs) της παραγράφου.
                    run_text = "".join(run.text or "" for run in paragraph.runs)
                    text_runs.append(run_text)
                
                # Κρατάω μόνο τις γραμμές που έχουν πραγματικό περιεχόμενο (όχι κενά).
                parts.append("\n".join([t for t in text_runs if t.strip()]))
        
        # Ενώνω όλα τα τμήματα κειμένου που βρέθηκαν στη διαφάνεια σε ένα ενιαίο string.
        slide_text = "\n".join([p for p in parts if p.strip()])
        
        # Αποθηκεύω το αποτέλεσμα (Αριθμός Διαφάνειας, Κείμενο) στη λίστα.
        results.append((i, slide_text))
        
    return results